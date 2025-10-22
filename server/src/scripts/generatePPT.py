#!/usr/bin/env python3
"""
generatePPT.py - Diagnostic and Fixed Version

Parses summary and populates PPTX template slides with bullet content.
Includes detailed logging to help identify placeholder issues.

Requires: py -3 -m pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import os
import re
import sys

NUM_SLIDES = 5
CANONICAL_HEADINGS = [
    "Background & Motivation",
    "Key Findings",
    "Methods & Evidence",
    "Therapeutic Implications",
    "Conclusion"
]


def load_payload(path):
    with open(path, "r", encoding="utf8") as fh:
        return json.load(fh)


def parse_summary_into_sections(summary, headings_list):
    """
    Parse a summary string into sections keyed by heading.
    - Ignores top-level "Detailed Summary with Key Points" heading.
    - Returns dict: heading -> content (list of bullet strings).
    """
    if not summary or not summary.strip():
        print("WARNING: Summary is empty", file=sys.stderr)
        return {h: [] for h in headings_list}

    # Normalize line endings
    txt = summary.replace("\r\n", "\n").replace("\r", "\n")
    
    print(f"DEBUG: Original summary length: {len(txt)}", file=sys.stderr)

    # Remove top-level label if present
    txt = re.sub(r'^\s*Detailed Summary with Key Points\s*[\n:]*', '', txt, flags=re.IGNORECASE | re.MULTILINE)

    # Build regex to find heading lines
    esc_headings = [re.escape(h.strip()) for h in headings_list]
    esc_headings.sort(key=len, reverse=True)
    pattern = r'(?im)^\s*(' + r'|'.join(esc_headings) + r')\s*$'

    lines = txt.splitlines()
    positions = []
    for idx, line in enumerate(lines):
        match = re.match(pattern, line.strip())
        if match:
            positions.append((idx, match.group(1).strip()))
            print(f"DEBUG: Found heading '{match.group(1).strip()}' at line {idx}", file=sys.stderr)

    sections = {h: [] for h in headings_list}
    
    if not positions:
        print("WARNING: No headings found in summary. Attempting fallback parsing.", file=sys.stderr)
        # Fallback: try to find headings with less strict pattern
        for idx, line in enumerate(lines):
            line_clean = line.strip()
            for h in headings_list:
                if h.lower() in line_clean.lower() and len(line_clean) < 100:
                    positions.append((idx, h))
                    print(f"DEBUG: Fallback found heading '{h}' at line {idx}", file=sys.stderr)
                    break
    
    if not positions:
        print("ERROR: Still no headings found after fallback", file=sys.stderr)
        return sections

    # Extract content between each heading and the next
    for i, (pos, hdr_line) in enumerate(positions):
        # Match heading to canonical (case-insensitive)
        matched = None
        for ch in headings_list:
            if hdr_line.lower() == ch.lower() or ch.lower() in hdr_line.lower():
                matched = ch
                break
        if not matched:
            print(f"WARNING: Could not match heading '{hdr_line}'", file=sys.stderr)
            continue

        start = pos + 1
        end = positions[i + 1][0] if i + 1 < len(positions) else len(lines)
        content_lines = lines[start:end]

        print(f"DEBUG: Extracting content for '{matched}' from line {start} to {end}", file=sys.stderr)

        # Parse bullet points
        bullets = []
        for ln in content_lines:
            stripped = ln.strip()
            if not stripped:
                continue
            # Remove leading bullet markers and whitespace
            cleaned = re.sub(r'^\s*[-•*]\s*', '', ln).strip()
            if cleaned:
                bullets.append(cleaned)

        print(f"DEBUG: Found {len(bullets)} bullets for '{matched}'", file=sys.stderr)
        sections[matched] = bullets

    return sections


def find_content_placeholder(slide, slide_title=""):
    """
    Find the content placeholder on the slide.
    Try multiple strategies:
    1. Look for placeholder with idx > 0 (not title)
    2. Look for any text box that's not the title
    3. Return position info for manual textbox creation
    """
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        pass

    print(f"DEBUG: Searching for content placeholder on slide '{slide_title}'", file=sys.stderr)
    print(f"DEBUG: Slide has {len(slide.shapes)} total shapes", file=sys.stderr)
    
    # Strategy 1: Try placeholders with specific indices
    try:
        placeholder_count = 0
        for shape in slide.placeholders:
            placeholder_count += 1
            try:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                print(f"DEBUG: Placeholder {placeholder_count}: idx={ph_idx}, type={ph_type}, has_text_frame={hasattr(shape, 'text_frame')}", file=sys.stderr)
                
                # Skip title placeholder (usually idx 0 or type 1)
                if shape == title_shape:
                    print(f"DEBUG: Skipping title placeholder", file=sys.stderr)
                    continue
                
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    print(f"DEBUG: Using placeholder idx={ph_idx} for content", file=sys.stderr)
                    return shape
            except Exception as e:
                print(f"DEBUG: Error checking placeholder {placeholder_count}: {e}", file=sys.stderr)
    except Exception as e:
        print(f"DEBUG: Error iterating placeholders: {e}", file=sys.stderr)

    # Strategy 2: Try all shapes (not just placeholders)
    print(f"DEBUG: Fallback to searching all shapes", file=sys.stderr)
    for idx, shape in enumerate(slide.shapes):
        if shape == title_shape:
            continue
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            print(f"DEBUG: Using shape {idx} for content", file=sys.stderr)
            return shape

    print(f"WARNING: No content placeholder found for slide '{slide_title}'", file=sys.stderr)
    return None


def clear_text_frame(tf):
    """Clear all text from a text frame."""
    while len(tf.paragraphs) > 1:
        try:
            p = tf.paragraphs[-1]
            p._element.getparent().remove(p._element)
        except Exception:
            break
    if tf.paragraphs:
        tf.paragraphs[0].text = ""


def set_bullets_into_shape(shape, bullets):
    """
    Write a list of bullet strings into a shape's text_frame as bullet points.
    """
    if not hasattr(shape, "text_frame"):
        print("ERROR: Shape has no text_frame", file=sys.stderr)
        return False
    
    tf = shape.text_frame
    clear_text_frame(tf)
    
    if not bullets:
        print("DEBUG: No bullets to write", file=sys.stderr)
        tf.text = ""
        return True

    print(f"DEBUG: Writing {len(bullets)} bullets to shape", file=sys.stderr)

    # Set first bullet
    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.level = 0
    try:
        p.font.size = Pt(14)
    except Exception:
        pass

    # Add remaining bullets
    for i, bullet_text in enumerate(bullets[1:], 1):
        p = tf.add_paragraph()
        p.text = bullet_text
        p.level = 0
        try:
            p.font.size = Pt(14)
        except Exception:
            pass
        print(f"DEBUG: Added bullet {i}: {bullet_text[:50]}...", file=sys.stderr)

    return True


def main():
    if len(sys.argv) < 2:
        print("Usage: generatePPT.py /path/to/payload.json", file=sys.stderr)
        sys.exit(2)
    
    payload_path = sys.argv[1]
    if not os.path.exists(payload_path):
        print("Payload not found: " + payload_path, file=sys.stderr)
        sys.exit(2)

    print(f"DEBUG: Loading payload from {payload_path}", file=sys.stderr)
    payload = load_payload(payload_path)
    
    template = payload.get("templatePath")
    output = payload.get("outputPath")
    summary = payload.get("summary", "")

    print(f"DEBUG: Template: {template}", file=sys.stderr)
    print(f"DEBUG: Output: {output}", file=sys.stderr)
    print(f"DEBUG: Summary length: {len(summary)}", file=sys.stderr)

    if not template or not os.path.exists(template):
        print("Template not found: {}".format(template), file=sys.stderr)
        sys.exit(2)

    # Parse summary into sections
    sections = parse_summary_into_sections(summary, CANONICAL_HEADINGS)
    
    # Print what we found
    for heading, bullets in sections.items():
        print(f"DEBUG: Section '{heading}' has {len(bullets)} bullets", file=sys.stderr)

    prs = Presentation(template)
    slides_list = list(prs.slides)
    print(f"DEBUG: Template has {len(slides_list)} slides", file=sys.stderr)

    # Build map of slide title (normalized) -> slide
    slide_map = {}
    for idx, slide in enumerate(slides_list):
        try:
            title_shape = slide.shapes.title
            if title_shape and hasattr(title_shape, "text"):
                title_text = title_shape.text.strip()
                if title_text:
                    slide_map[title_text.lower()] = slide
                    print(f"DEBUG: Slide {idx} title: '{title_text}'", file=sys.stderr)
        except Exception as e:
            print(f"DEBUG: Slide {idx} has no title: {e}", file=sys.stderr)

    # Match each canonical heading to a slide and populate content
    for heading in CANONICAL_HEADINGS:
        bullets = sections.get(heading, [])
        
        print(f"\nDEBUG: Processing heading '{heading}' with {len(bullets)} bullets", file=sys.stderr)
        
        # Find slide by matching title (case-insensitive)
        target_slide = None
        heading_lower = heading.lower()
        
        # Exact match
        if heading_lower in slide_map:
            target_slide = slide_map[heading_lower]
            print(f"DEBUG: Found exact match for '{heading}'", file=sys.stderr)
        else:
            # Fuzzy match
            for slide_title_lower, slide in slide_map.items():
                if heading_lower in slide_title_lower or slide_title_lower in heading_lower:
                    target_slide = slide
                    print(f"DEBUG: Found fuzzy match for '{heading}' -> '{slide_title_lower}'", file=sys.stderr)
                    break

        if not target_slide:
            print(f"ERROR: No slide found for heading '{heading}'", file=sys.stderr)
            continue

        # Find content placeholder and write bullets
        content_shape = find_content_placeholder(target_slide, heading)
        if content_shape:
            success = set_bullets_into_shape(content_shape, bullets)
            if success:
                print(f"SUCCESS: Wrote {len(bullets)} bullets to slide '{heading}'", file=sys.stderr)
        else:
            # No placeholder found; add a textbox as fallback
            print(f"WARNING: Adding fallback textbox for slide '{heading}'", file=sys.stderr)
            tbox = target_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.0))
            set_bullets_into_shape(tbox, bullets)

    # Save
    try:
        prs.save(output)
        print("OK: saved to " + output)
        sys.exit(0)
    except Exception as e:
        print("Failed to save PPTX: " + str(e), file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()